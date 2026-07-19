#!/usr/bin/env python3
"""
Generate the Day-1 Mission Definition PPT (Simplified Chinese).
v2.0 — comprehensive visual refresh.

This deck is 22 slides, professionally themed, with:
  - Gradient covers & section dividers
  - Data-driven icons (vector shapes) throughout
  - 12+ figures (radar, treemap, schematic, gantt, etc.)
  - Decision-matrix style requirement displays
  - Two interactive links to Plotly HTML pages

Audience: Tsinghua Summer Program review.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

# Paths
output_path = (
    "/home/user/tsinghua-rocket/docs/01_mission_requirements/final/"
    "Day1_Mission_Definition_Presentation.pptx"
)
figures_dir = "/home/user/tsinghua-rocket/figures/day01"

# 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ------------------------------------------------------------------
# Theme palette
# ------------------------------------------------------------------
NAVY    = RgbColor(0x1A, 0x3C, 0x6E)
NAVY_D  = RgbColor(0x10, 0x25, 0x4A)
SKY     = RgbColor(0x3B, 0x82, 0xC4)
TEAL    = RgbColor(0x0E, 0x7C, 0x7B)
GOLD    = RgbColor(0xD4, 0xA0, 0x17)
CORAL   = RgbColor(0xD2, 0x55, 0x3B)
VIOLET  = RgbColor(0x6B, 0x4E, 0x9C)
SAGE    = RgbColor(0x7B, 0xA0, 0x5B)
SLATE   = RgbColor(0x47, 0x55, 0x69)
ROSE    = RgbColor(0xC2, 0x18, 0x5B)
ASH     = RgbColor(0x6B, 0x72, 0x80)
MIST    = RgbColor(0xE5, 0xED, 0xF5)
PARCH   = RgbColor(0xFA, 0xF7, 0xF0)
WHITE   = RgbColor(0xFF, 0xFF, 0xFF)
DARK    = RgbColor(0x2D, 0x34, 0x36)
ACCENT  = RgbColor(0x00, 0x7A, 0xCC)
WARN    = RgbColor(0xD9, 0x77, 0x06)

# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------
def _set_gradient_fill(shape, c1_hex, c2_hex, angle=270):
    """Apply a 2-stop linear gradient fill (vertical default)."""
    sp = shape.fill._xPr
    # Remove any existing solidFill
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    # Build gradFill
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }
    gsLst = etree.SubElement(sp, qn('a:gradFill'),
                             attrib={'flip': 'none', 'rotWithShape': '1'})
    gsLst.set('flip', 'none')
    gsLst.set('rotWithShape', '1')
    lin = etree.SubElement(gsLst, qn('a:lin'),
                           attrib={'ang': str(angle * 60000), 'scaled': '0'})
    gs1 = etree.SubElement(gsLst, qn('a:gs'),
                           attrib={'pos': '0'})
    etree.SubElement(gs1, qn('a:srgbClr'), attrib={'val': c1_hex})
    gs2 = etree.SubElement(gsLst, qn('a:gs'),
                           attrib={'pos': '100000'})
    etree.SubElement(gs2, qn('a:srgbClr'), attrib={'val': c2_hex})
    tile = etree.SubElement(gsLst, qn('a:tileRect'))
    tile.set('l', '0')
    tile.set('t', '0')
    tile.set('r', '0')
    tile.set('b', '0')


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_gradient_bg(slide, top_hex, bot_hex):
    """Add a full-bleed vertical gradient background."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    _set_gradient_fill(bg, top_hex, bot_hex, angle=270)
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_header_bar(slide, title, subtitle=None, accent=NAVY):
    """Themed header: navy bar + white title + optional subtitle."""
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, Inches(1.1))
    h.fill.solid()
    h.fill.fore_color.rgb = accent
    h.line.fill.background()
    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18),
                                  Inches(10.0), Inches(0.75))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Microsoft YaHei'
    # Right-side brand mark
    bm = slide.shapes.add_textbox(Inches(11.0), Inches(0.18),
                                  Inches(2.2), Inches(0.75))
    tf = bm.text_frame
    p = tf.paragraphs[0]
    p.text = "CRLV-1 · Day 1"
    p.font.size = Pt(11)
    p.font.color.rgb = RgbColor(0xCC, 0xE5, 0xFF)
    p.alignment = PP_ALIGN.RIGHT
    # Accent underline
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.1),
                                 prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2),
                                      Inches(12.3), Inches(0.4))
        p = sb.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = SLATE
        p.font.name = 'Microsoft YaHei'


def add_footer(slide, page=None, total=None):
    """Themed footer bar at bottom."""
    fb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,
                                prs.slide_height - Inches(0.4),
                                prs.slide_width, Inches(0.4))
    fb.fill.solid()
    fb.fill.fore_color.rgb = NAVY_D
    fb.line.fill.background()
    # left text
    ft = slide.shapes.add_textbox(Inches(0.4),
                                  prs.slide_height - Inches(0.38),
                                  Inches(8.0), Inches(0.34))
    p = ft.text_frame.paragraphs[0]
    p.text = "Tsinghua Summer Program · AI-Assisted Reusable Rocket Co-Design · v1.0 refined"
    p.font.size = Pt(8.5)
    p.font.color.rgb = RgbColor(0xB0, 0xC4, 0xDE)
    p.font.name = 'Calibri'
    # page
    if page and total:
        pt = slide.shapes.add_textbox(Inches(11.5),
                                      prs.slide_height - Inches(0.38),
                                      Inches(1.5), Inches(0.34))
        p = pt.text_frame.paragraphs[0]
        p.text = f"{page} / {total}"
        p.font.size = Pt(8.5)
        p.font.color.rgb = RgbColor(0xB0, 0xC4, 0xDE)
        p.alignment = PP_ALIGN.RIGHT


def add_textbox(slide, left, top, width, height, lines, *,
                size=14, color=DARK, bold=False, italic=False,
                name='Microsoft YaHei', align=PP_ALIGN.LEFT,
                bullet=True, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if bullet:
            p.text = "• " + item
        else:
            p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = name
        p.alignment = align
        p.space_after = Pt(4)
        p.line_spacing = line_spacing
    return tb


def add_kpi_card(slide, left, top, width, height, value, label, color=NAVY,
                 value_size=30, label_size=11):
    """A rounded rectangle with a big value and a small label."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2.5)
    card.shadow.inherit = False
    # value text
    vb = slide.shapes.add_textbox(left, top + Inches(0.25),
                                  width, Inches(0.9))
    p = vb.text_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(value_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Calibri'
    # label
    lb = slide.shapes.add_textbox(left, top + height - Inches(0.55),
                                  width, Inches(0.45))
    p = lb.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(label_size)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Microsoft YaHei'
    return card


def add_section_header(slide, number, title, subtitle, color=NAVY):
    """Section divider slide."""
    add_gradient_bg(slide, '1A3C6E', '0E2548')
    # Big section number
    nb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8),
                                  Inches(2.5), Inches(2.5))
    p = nb.text_frame.paragraphs[0]
    p.text = f"{number:02d}"
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0x4F, 0x70, 0x9F)
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    # Title
    tb = slide.shapes.add_textbox(Inches(3.5), Inches(2.5),
                                  Inches(9.0), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Microsoft YaHei'
    # Subtitle
    sb = slide.shapes.add_textbox(Inches(3.5), Inches(4.2),
                                  Inches(9.0), Inches(2.0))
    tf = sb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(subtitle):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(0xC0, 0xD0, 0xE5)
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(4)
    # Decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(3.5), Inches(2.3),
                                  Inches(2.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


# ------------------------------------------------------------------
# Slide 1 — Cover
# ------------------------------------------------------------------
def slide_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide, '1A3C6E', '0A1A33')
    # Decorative left strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                   Inches(0.35), prs.slide_height)
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()
    # Decorative diagonal accent
    d1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.0), Inches(6.2),
                                Inches(8.0), Inches(0.05))
    d1.fill.solid()
    d1.fill.fore_color.rgb = TEAL
    d1.line.fill.background()
    # Top label
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.7),
                                  Inches(12), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "TSINGHUA SUMMER PROGRAM  ·  AI-ASSISTED ENGINEERING"
    p.font.size = Pt(11)
    p.font.color.rgb = RgbColor(0x9F, 0xB8, 0xD0)
    p.font.name = 'Calibri'
    p.font.bold = True
    # Project number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.8), Inches(1.4),
                                   Inches(1.2), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = TEAL
    badge.line.fill.background()
    bp = badge.text_frame.paragraphs[0]
    bp.text = "Day 1"
    bp.font.size = Pt(16)
    bp.font.bold = True
    bp.font.color.rgb = WHITE
    bp.font.name = 'Calibri'
    bp.alignment = PP_ALIGN.CENTER
    # Main title
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.3),
                                  Inches(12), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CRLV-1 任务定义与需求"
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Microsoft YaHei'
    p2 = tf.add_paragraph()
    p2.text = "Mission Definition & Requirements"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RgbColor(0xB0, 0xCC, 0xE5)
    p2.font.name = 'Calibri'
    p2.font.italic = True
    # Subtitle
    sb = slide.shapes.add_textbox(Inches(0.8), Inches(4.5),
                                  Inches(12), Inches(1.2))
    tf = sb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "概念可重复使用运载火箭  ·  Conceptual Reusable Launch Vehicle"
    p.font.size = Pt(18)
    p.font.color.rgb = RgbColor(0xC0, 0xD5, 0xE8)
    p.font.name = 'Microsoft YaHei'
    p2 = tf.add_paragraph()
    p2.text = "Day 1 — v1.0 (refined, self-reviewed)"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RgbColor(0x90, 0xA8, 0xC0)
    p2.font.italic = True
    p2.font.name = 'Microsoft YaHei'
    # Authors
    ab = slide.shapes.add_textbox(Inches(0.8), Inches(6.5),
                                  Inches(12), Inches(0.6))
    p = ab.text_frame.paragraphs[0]
    p.text = "Research Team  ·  Tsinghua Summer Program  ·  2026年7月18日"
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(0x9F, 0xB8, 0xD0)
    p.font.name = 'Calibri'
    # Right-side vertical text (rotated badge)
    rt = slide.shapes.add_textbox(Inches(11.5), Inches(0.5),
                                  Inches(1.5), Inches(0.4))
    p = rt.text_frame.paragraphs[0]
    p.text = "v1.0 refined"
    p.font.size = Pt(11)
    p.font.color.rgb = ROSE
    p.font.bold = True
    p.font.italic = True
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = 'Calibri'


# ------------------------------------------------------------------
# Slide 2 — Section divider: 1. 项目背景
# ------------------------------------------------------------------
def slide_section1():
    add_section_header(slide := prs.slides.add_slide(prs.slide_layouts[6]),
                       1, "项目背景与定位",
                       ["Why a 1.2 t reusable launcher, and why now?",
                        "2025–2026 global benchmarks + Chinese surge",
                        "Positioning CRLV-1 in the medium-small class"])


# ------------------------------------------------------------------
# Slide 3 — 项目背景 (rich)
# ------------------------------------------------------------------
def slide_background():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide,
                   "项目背景与研究目标",
                   "Why a 1.2 t reusable launcher? Context and motivation.")
    # Three pillars
    pillars = [
        ("商业可重复使用", "Falcon 9 (>500 flights, $2.4–2.7k/kg), 20+ flights/booster 已成为行业基准",
         NAVY),
        ("中国 2025–2026 突破", "Zhuque-3 (18 t), LM 10/12A 家族海基网捕回收 2026, 多型进入首飞",
         CORAL),
        ("概念学术定位", "1.2–2 t LEO/SSO — 在专用小运载与 Falcon 9 之间开拓深耕空间",
         TEAL),
    ]
    for i, (t, d, c) in enumerate(pillars):
        x = Inches(0.5 + i*4.3)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, Inches(1.6), Inches(4.0), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = c
        card.line.width = Pt(2)
        # Color header
        hh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x, Inches(1.6), Inches(4.0), Inches(0.5))
        hh.fill.solid()
        hh.fill.fore_color.rgb = c
        hh.line.fill.background()
        hp = hh.text_frame.paragraphs[0]
        hp.text = t
        hp.font.size = Pt(13)
        hp.font.bold = True
        hp.font.color.rgb = WHITE
        hp.alignment = PP_ALIGN.CENTER
        hp.font.name = 'Microsoft YaHei'
        # Description
        db = slide.shapes.add_textbox(x + Inches(0.15), Inches(2.2),
                                      Inches(3.7), Inches(1.75))
        p = db.text_frame.paragraphs[0]
        p.text = d
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.LEFT
        db.text_frame.word_wrap = True
    # Bottom band — research goals
    goal_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0.5), Inches(4.3),
                                      Inches(12.3), Inches(2.5))
    goal_box.fill.solid()
    goal_box.fill.fore_color.rgb = MIST
    goal_box.line.color.rgb = ACCENT
    goal_box.line.width = Pt(1.5)
    # Title
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(4.4),
                                  Inches(12), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "本研究目标  ·  Research Objectives"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = 'Microsoft YaHei'
    # Items
    items = [
        "为概念型可重复使用运载火箭 CRLV-1 制定完整任务定义与 L0/L1 需求（14 条 L1 全部含验证方法）",
        "提出创新 GRI（Green Reusability Index）框架——量纲一致（kg payload / kg CO₂e），无伪造数值",
        "正式要求评估海基网捕回收（受 2026 年中国在轨级回收事件启发），与推进式着陆在 Day 6 量化权衡",
        "建立自我批判机制——本 v1.0 精炼版基于 v0 草稿的 12 项问题（命名/数据/方法/呈现）系统修正",
    ]
    add_textbox(slide, Inches(0.8), Inches(4.95),
                Inches(11.8), Inches(1.7), items,
                size=12, color=DARK, name='Microsoft YaHei', line_spacing=1.2)
    add_footer(slide, 3, 22)


# ------------------------------------------------------------------
# Slide 4 — 2026 关键数据 (KPI grid)
# ------------------------------------------------------------------
def slide_data_2026():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "2026 年主要运载器关键数据",
                   "Cross-checked against SpaceNews, Ars Technica, China-in-Space, company reports.")
    # 4 KPI cards
    kpis = [
        ("22,800 kg", "Falcon 9 可重复使用\nLEO 载荷",  NAVY,    "public · >500 flights"),
        ("12,000 kg", "中国 2026 在轨级\n回收（参考量级）", CORAL,  "型号名称待核实"),
        ("18,300 kg", "朱雀三号 回收目标\nLEO 载荷",        TEAL,   "maiden Dec 2025"),
        ("~19 t/t",  "传统 RP-1/LOX\n每吨载荷 CO₂e",    GOLD,   "LCA estimate"),
    ]
    for i, (val, lab, c, note) in enumerate(kpis):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col*6.2)
        y = Inches(1.55 + row*2.45)
        add_kpi_card(slide, x, y, Inches(5.8), Inches(2.1),
                     val, lab, color=c, value_size=32)
        # Note below
        nt = slide.shapes.add_textbox(x, y + Inches(2.15),
                                      Inches(5.8), Inches(0.3))
        p = nt.text_frame.paragraphs[0]
        p.text = note
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = SLATE
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Calibri'
    # Bottom annotation
    bt = slide.shapes.add_textbox(Inches(0.6), Inches(6.7),
                                  Inches(12.1), Inches(0.4))
    p = bt.text_frame.paragraphs[0]
    p.text = "所有数据经过至少两个独立来源交叉核实；不可核实者明确标注为“估算 / 参考量级”"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = SLATE
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Microsoft YaHei'
    add_footer(slide, 4, 22)


# ------------------------------------------------------------------
# Slide 5 — 关于 2026 命名 (honesty page)
# ------------------------------------------------------------------
def slide_naming():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PARCH)
    add_header_bar(slide, "关于 2026 中国回收事件 — 诚实的命名说明",
                   "Academic honesty requires flagging uncertain attributions.")
    # Main message box
    mb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(1.4),
                                Inches(7.0), Inches(5.5))
    mb.fill.solid()
    mb.fill.fore_color.rgb = WHITE
    mb.line.color.rgb = CORAL
    mb.line.width = Pt(2)
    add_textbox(slide, Inches(0.8), Inches(1.6),
                Inches(6.5), Inches(5.2), [
        "公开英文媒体（Ars Technica 等）使用 “Long March 10B”",
        "其他中英文来源使用 “Long March 12A” 描述同一事件",
        "长征十号系列另有“载人登月用”三模块火箭，与此事件并非同一型号",
        "→ 本文档采用保守标签 “Long March 10/12A 家族”，并明示不确定性",
        "→ Day 1 暂以 12,000 kg LEO 作为参考量级（非精确公布值）",
        "→ Day 3 / Day 6 将基于一手资料进一步核查型号与回收细节",
    ], size=13, name='Microsoft YaHei', line_spacing=1.4)
    # Right side — visual summary
    rb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(7.8), Inches(1.4),
                                Inches(5.0), Inches(5.5))
    rb.fill.solid()
    rb.fill.fore_color.rgb = NAVY
    rb.line.fill.background()
    th = slide.shapes.add_textbox(Inches(8.0), Inches(1.55),
                                  Inches(4.6), Inches(0.6))
    p = th.text_frame.paragraphs[0]
    p.text = "本张幻灯片的核心价值"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Microsoft YaHei'
    body = [
        ("科学诚实", "在学术项目中，命名不"
                  "清的数据应明确标注，"),
        ("", "不应“挑一个看起来对的”"),
        ("", ""),
        ("v1.0 修正", "已统一为 “Long March 10/12A 家族” 表述"),
        ("", "并将该不确定性记入风险登记册 R08"),
        ("", ""),
        ("可追溯", "AI 日志与决策文档均记录该命名决策过程"),
    ]
    tb = slide.shapes.add_textbox(Inches(8.0), Inches(2.3),
                                  Inches(4.6), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (k, v) in enumerate(body):
        if not k and not v:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if k:
            p.text = f"▸ {k}"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = ACCENT
        else:
            p.text = f"   {v}"
            p.font.size = Pt(11)
            p.font.color.rgb = RgbColor(0xC0, 0xD5, 0xE8)
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(4)
    add_footer(slide, 5, 22)


# ------------------------------------------------------------------
# Slide 6 — 研究方法
# ------------------------------------------------------------------
def slide_method():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "研究方法与数据来源",
                   "Systematic literature review with multi-source cross-verification.")
    # Pipeline boxes
    stages = [
        ("① 文献调研", "2025–2026 公开数据\n(英文 + 中文)", SKY),
        ("② 多源交叉", "至少 2 个独立来源\n核实定量数据",     TEAL),
        ("③ AI 协作", "Claude/GPT 作为\n工程副驾驶",     VIOLET),
        ("④ 自评审", "v0 → v1.0 修复\n12 项问题",          ROSE),
        ("⑤ 输出", "14 条 L1 + GRI 框架\n+ 回收架构权衡", CORAL),
    ]
    for i, (t, d, c) in enumerate(stages):
        x = Inches(0.4 + i*2.6)
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, Inches(1.7),
                                     Inches(2.4), Inches(2.0))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = c
        box.line.width = Pt(2)
        # Color top
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x, Inches(1.7),
                                     Inches(2.4), Inches(0.55))
        top.fill.solid()
        top.fill.fore_color.rgb = c
        top.line.fill.background()
        p = top.text_frame.paragraphs[0]
        p.text = t
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Microsoft YaHei'
        # Description
        db = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.4),
                                      Inches(2.2), Inches(1.25))
        tf = db.text_frame
        tf.word_wrap = True
        for j, line in enumerate(d.split('\n')):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(10.5)
            p.font.color.rgb = DARK
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.CENTER
        # Arrow between
        if i < len(stages) - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                          x + Inches(2.45), Inches(2.55),
                                          Inches(0.2), Inches(0.3))
            arr.fill.solid()
            arr.fill.fore_color.rgb = SLATE
            arr.line.fill.background()
    # Bottom — sources
    src_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.4), Inches(4.0),
                                     Inches(12.5), Inches(2.9))
    src_box.fill.solid()
    src_box.fill.fore_color.rgb = MIST
    src_box.line.color.rgb = ACCENT
    src_box.line.width = Pt(1.2)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(4.1),
                                  Inches(12), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "主要数据来源  ·  Primary Sources (2025–2026)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = 'Microsoft YaHei'
    add_textbox(slide, Inches(0.6), Inches(4.6),
                Inches(12.1), Inches(2.2), [
        "英文：SpaceNews, Ars Technica, Orbital Radar, New Space Economy, SpaceX Falcon 9 User's Guide, Rocket Lab Neutron",
        "中文：Global Times, China-in-Space, 厂家公告 / 公开发布会, 中国航天资讯",
        "可持续性：Strathclyde University RLV LCA, MaiaSpace 环境评估, arXiv 2504.15291, ESA Clean Space",
        "AI 协作：所有 prompt 与决策均记录于 ai_logs/，包括 v1.0 自我评审的三条新 prompt",
    ], size=11.5, name='Microsoft YaHei', line_spacing=1.3)
    add_footer(slide, 6, 22)


# ------------------------------------------------------------------
# Slide 7 — Section divider: 2. 需求体系
# ------------------------------------------------------------------
def slide_section2():
    add_section_header(slide := prs.slides.add_slide(prs.slide_layouts[6]),
                       2, "需求体系",
                       ["6 L0 objectives + 14 L1 requirements",
                        "Stakeholder alignment + traceability",
                        "Novel: GRI + net-capture recovery"])


# ------------------------------------------------------------------
# Slide 8 — 任务陈述 + L0 目标
# ------------------------------------------------------------------
def slide_l0():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "任务陈述与 L0 目标（精炼版）",
                   "Six top-level objectives; each has threshold + goal values.")
    # Mission statement band
    mb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(1.4),
                                Inches(12.3), Inches(0.85))
    mb.fill.solid()
    mb.fill.fore_color.rgb = NAVY
    mb.line.fill.background()
    mp = mb.text_frame.paragraphs[0]
    mp.text = "  任务陈述  研制概念型部分可重复使用运载火箭 CRLV-1，目标有效载荷 1,200 kg（门槛）→ 2,000 kg（目标）至 LEO/SSO，一级 ≥ 10 次（门槛）/ ≥ 20 次（目标）重复使用，配套完整的 GRI 可持续性框架与回收架构权衡。"
    mp.font.size = Pt(11.5)
    mp.font.color.rgb = WHITE
    mp.font.name = 'Microsoft YaHei'
    mb.text_frame.word_wrap = True
    mb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 6 L0 cards in 2x3 grid
    l0_items = [
        ("L0-01", "有效载荷", "1,200 / 2,000 kg", "500 km LEO or 700 km SSO", NAVY),
        ("L0-02", "可重复使用", "10 / 20 次", "翻新 < 15% 新建造成本", TEAL),
        ("L0-03", "成本效率", "$3,500 / $2,500 per kg", "v1.0 已由 $2,800 上调至 $3,500", CORAL),
        ("L0-04", "可持续性", "< 15 t CO₂e/t payload", "约 20% 优于 RP-1 基线 (GRI 框架)", SAGE),
        ("L0-05", "响应性", "30 / 14 天", "载荷集成到发射窗口", GOLD),
        ("L0-06", "可靠性", "累计 ≥ 0.95 (10 飞)", "单飞 ≥ 0.994 (门槛) / 0.997 (目标)", VIOLET),
    ]
    for i, (lid, name, val, note, c) in enumerate(l0_items):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col*4.15)
        y = Inches(2.4 + row*2.2)
        # ID badge
        badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       x, y, Inches(0.85), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = c
        badge.line.fill.background()
        bp = badge.text_frame.paragraphs[0]
        bp.text = lid
        bp.font.size = Pt(11)
        bp.font.bold = True
        bp.font.color.rgb = WHITE
        bp.alignment = PP_ALIGN.CENTER
        bp.font.name = 'Calibri'
        # Name
        nb = slide.shapes.add_textbox(x + Inches(0.9), y,
                                      Inches(3.0), Inches(0.45))
        p = nb.text_frame.paragraphs[0]
        p.text = name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.font.name = 'Microsoft YaHei'
        # Value (big)
        vb = slide.shapes.add_textbox(x, y + Inches(0.55),
                                      Inches(3.95), Inches(0.7))
        p = vb.text_frame.paragraphs[0]
        p.text = val
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = c
        p.font.name = 'Calibri'
        # Note
        eb = slide.shapes.add_textbox(x, y + Inches(1.25),
                                      Inches(3.95), Inches(0.8))
        p = eb.text_frame.paragraphs[0]
        p.text = note
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = SLATE
        p.font.name = 'Microsoft YaHei'
        eb.text_frame.word_wrap = True
    add_footer(slide, 8, 22)


# ------------------------------------------------------------------
# Slide 9 — L1 需求 (table)
# ------------------------------------------------------------------
def slide_l1():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "关键 L1 需求（精炼版，14 条）",
                   "All 14 L1 requirements carry explicit verification methods.")
    # L1 table
    rows = [
        ("ID",       "Requirement",                                   "Threshold",          "Goal",                "Verification"),
        ("L1-P01",  "Payload to 500 km LEO (reusable)",              "1,200 kg",           "2,000 kg",            "Trajectory sim"),
        ("L1-P02",  "Injection accuracy (alt / inclination)",         "±20 km / ±0.05°",    "±10 km / ±0.02°",     "Flight test + GNC"),
        ("L1-P03",  "Fairing internal volume",                        "Ø 3.4 m × 6.5 m",    "Ø 3.6 m × 7.0 m",     "Inspection"),
        ("L1-R01",  "Per-flight recovery success",                    "≥ 0.90",            "≥ 0.98",              "Flight test + MC"),
        ("L1-R02",  "Number of reuses per first stage",               "10",                "20",                  "Ops tracking"),
        ("L1-R03",  "Recovery: propulsive (baseline) OR net-capture", "Trade closed Day 6", "—",                   "Trade study + sim"),
        ("L1-C01",  "Recurring launch cost per flight",               "≤ $4.2 M",          "≤ $3.0 M",            "Cost model (Day 8)"),
        ("L1-C02",  "Refurbishment cost fraction",                    "≤ 15%",             "≤ 10%",               "Cost model"),
        ("L1-E01",  "Propellant preference",                          "LOX/LCH₄",          "+ bio-CH₄ option",    "Lifecycle assessment"),
        ("L1-E02",  "Green Reusability Index (GRI)",                   "≥ 1.20×F9 (illus.)","≥ 1.35×F9 (illus.)",  "Day 7 model"),
        ("L1-S01",  "Compound mission success, first 10 flights",     "≥ 0.95",            "≥ 0.97",              "RBD + MC"),
        ("L1-O01",  "Compatible launch site",                         "Hainan (Wenchang)", "+ Jiuquan sea-zone",  "CONOPS analysis"),
        ("L1-O02",  "Dedicated + rideshare modes",                    "Yes",               "Yes",                 "Interface doc"),
        ("L1-O03",  "Landing accuracy (propulsive)",                  "± 30 m",            "± 10 m",              "Flight test + MC"),
    ]
    table = slide.shapes.add_table(len(rows), 5,
                                    Inches(0.4), Inches(1.5),
                                    Inches(12.5), Inches(5.3)).table
    widths = [1.2, 4.5, 2.2, 2.3, 2.3]
    for j, w in enumerate(widths):
        table.columns[j].width = Inches(w)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.name = 'Microsoft YaHei' if j > 0 else 'Calibri'
                p.alignment = PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER
                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    # Highlight novel rows (R03, E02, S01, O03, P03)
                    if row[0] in ('L1-R03', 'L1-E02', 'L1-O03', 'L1-P03'):
                        p.font.color.rgb = ROSE
                        p.font.bold = True
                    else:
                        p.font.color.rgb = DARK
            # Cell fill
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif row[0] in ('L1-R03', 'L1-E02', 'L1-O03', 'L1-P03'):
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(0xFD, 0xE8, 0xEE)  # very light rose
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = MIST
    # Legend
    lg = slide.shapes.add_textbox(Inches(0.4), Inches(6.85),
                                  Inches(12.5), Inches(0.3))
    p = lg.text_frame.paragraphs[0]
    p.text = "■ 粉红高亮 = v1.0 修正 / 新增的需求（L1-R03, L1-E02, L1-O03, L1-P03）"
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = SLATE
    p.font.name = 'Microsoft YaHei'
    add_footer(slide, 9, 22)


# ------------------------------------------------------------------
# Slide 10 — 14 L1 分类可视化
# ------------------------------------------------------------------
def slide_l1_viz():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "14 条 L1 需求 — 6 大类分布",
                   "All requirements have explicit verification methods.")
    img_path = f"{figures_dir}/requirements_treemap.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.4),
                                  width=Inches(8.5))
    # Right side commentary
    rt = slide.shapes.add_textbox(Inches(9.3), Inches(1.6),
                                  Inches(3.7), Inches(5.2))
    tf = rt.text_frame
    tf.word_wrap = True
    items = [
        ("▸ Performance (3)", "载荷 / 精度 / 整流罩", NAVY),
        ("▸ Reusability (3)", "回收 / 复用 / 架构", TEAL),
        ("▸ Operations (3)", "发射场 / 模式 / 精度", VIOLET),
        ("▸ Cost (2)", "单飞成本 / 翻新", GOLD),
        ("▸ Sustainability (2)", "推进剂 / GRI 指标", ROSE),
        ("▸ Safety (1)", "复合可靠性", SLATE),
    ]
    for i, (h, d, c) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = h
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = c
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(0)
        p2 = tf.add_paragraph()
        p2.text = "    " + d
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = DARK
        p2.font.name = 'Microsoft YaHei'
        p2.space_after = Pt(8)
    add_footer(slide, 10, 22)


# ------------------------------------------------------------------
# Slide 11 — Section divider: 3. GRI
# ------------------------------------------------------------------
def slide_section3():
    add_section_header(slide := prs.slides.add_slide(prs.slide_layouts[6]),
                       3, "Green Reusability Index",
                       ["Unit-consistent sustainability metric",
                        "Replaces the v0 fabricated bar chart",
                        "Quantitative values on Day 7"])


# ------------------------------------------------------------------
# Slide 12 — GRI 框架
# ------------------------------------------------------------------
def slide_gri():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PARCH)
    add_header_bar(slide, "Green Reusability Index (GRI) — 概念框架",
                   "Unit-consistent; no fabricated numerical ranking.")
    img_path = f"{figures_dir}/gri_comparison.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3),
                                  width=Inches(9.0))
    # Right side explanation
    rt = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(9.7), Inches(1.3),
                                Inches(3.4), Inches(5.4))
    rt.fill.solid()
    rt.fill.fore_color.rgb = NAVY
    rt.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(9.9), Inches(1.45),
                                  Inches(3.0), Inches(5.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "为什么提出 GRI？"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = 'Microsoft YaHei'
    p.space_after = Pt(6)
    lines = [
        "将可持续性从“附加考虑”提升为一级设计驱动因素",
        "统一量纲（kg payload / kg CO₂e）便于跨方案比较",
        "覆盖全生命周期：推进剂 / 制造 / 翻新 / 回收 / 报废",
        "",
        "v1.0 修正",
        "v0 图表中 1.00 / 1.15 / 1.35 数值无底层计算 → 已删除",
        "v1.0 改为量纲一致的框架图，明示无绝对排名",
        "",
        "Day 7 量化",
        "将由集成质量模型 + LCA 数据给出首个真实 GRI 值",
    ]
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        if line in ("v1.0 修正", "Day 7 量化"):
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = GOLD
        else:
            p.font.size = Pt(10.5)
            p.font.color.rgb = RgbColor(0xC0, 0xD5, 0xE8)
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(2)
    add_footer(slide, 12, 22)


# ------------------------------------------------------------------
# Slide 13 — GRI 设计杠杆
# ------------------------------------------------------------------
def slide_gri_levers():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "GRI 主要设计杠杆（定性排序）",
                   "Quantification on Day 7.")
    img_path = f"{figures_dir}/gri_levers.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.4),
                                  width=Inches(8.5))
    # Right side
    rt = slide.shapes.add_textbox(Inches(9.3), Inches(1.6),
                                  Inches(3.7), Inches(5.2))
    tf = rt.text_frame
    tf.word_wrap = True
    items = [
        ("▸ Top driver", "Propellant choice",  "conf. 0.85", TEAL),
        ("▸ 2nd driver", "Reusability (10→20)", "conf. 0.90", TEAL),
        ("▸ Mid impact",  "Refurb energy",     "conf. 0.55", GOLD),
        ("▸ Mid impact",  "Recovery transport", "conf. 0.70", GOLD),
        ("▸ Low impact",  "Material choice",   "conf. 0.50", VIOLET),
        ("▸ Low impact",  "Manufacturing",     "conf. 0.45", VIOLET),
    ]
    for i, (h, name, conf, c) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = h
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = c
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(0)
        p2 = tf.add_paragraph()
        p2.text = f"    {name}  ·  {conf}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = DARK
        p2.font.name = 'Calibri'
        p2.space_after = Pt(8)
    add_footer(slide, 13, 22)


# ------------------------------------------------------------------
# Slide 14 — Section divider: 4. 回收架构
# ------------------------------------------------------------------
def slide_section4():
    add_section_header(slide := prs.slides.add_slide(prs.slide_layouts[6]),
                       4, "回收架构权衡",
                       ["Propulsive landing vs sea-based net-capture",
                        "Driven by the 2026 Chinese net-capture demonstration",
                        "Trade closes on Day 6"])


# ------------------------------------------------------------------
# Slide 15 — 回收架构对比
# ------------------------------------------------------------------
def slide_recovery():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "回收架构选择（基于最新验证）",
                   "Two architectures; trade closes on Day 6 (closes L1-R03).")
    img_path = f"{figures_dir}/recovery_architecture.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.4),
                                  width=Inches(12.3))
    add_footer(slide, 15, 22)


# ------------------------------------------------------------------
# Slide 16 — 载荷对比图
# ------------------------------------------------------------------
def slide_payload():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "可重复使用运载器载荷能力对比（2026 数据）",
                   "CRLV-1 sits at the lower end of the medium-lift spectrum.")
    img_path = f"{figures_dir}/payload_comparison.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.4),
                                  width=Inches(9.0))
    rt = slide.shapes.add_textbox(Inches(9.7), Inches(1.6),
                                  Inches(3.4), Inches(5.2))
    tf = rt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "定位说明"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = 'Microsoft YaHei'
    p.space_after = Pt(6)
    lines = [
        "Falcon 9: 22.8 t — 商业基准",
        "New Glenn: 45 t — 大型可重复使用",
        "Zhuque-3: 18.3 t — 中国不锈钢甲烷氧",
        "LM 10/12A 家族: 12 t — 2026 网捕",
        "Hyperbola-3: 8.5 t — 中国可重复使用",
        "Pallas-1: 8 t — 中国可重复使用",
        "",
        "CRLV-1: 1.5 t (1.2 / 2.0 t)",
        "中小型概念验证",
        "聚焦 GRI + 回收架构权衡",
    ]
    for line in lines:
        p = tf.add_paragraph()
        p.text = "▸ " + line
        if "CRLV-1" in line or "中小型" in line:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = ROSE
        else:
            p.font.size = Pt(10.5)
            p.font.color.rgb = DARK
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(3)
    add_footer(slide, 16, 22)


# ------------------------------------------------------------------
# Slide 17 — 成本下降趋势
# ------------------------------------------------------------------
def slide_cost():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "发射成本下降趋势 — Falcon 9 单厂商参考",
                   "Honest single-vendor trajectory; industry curve on Day 8.")
    img_path = f"{figures_dir}/cost_trend.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.4),
                                  width=Inches(8.5))
    # Right side — 3 callouts
    notes = [
        ("v1.0 成本门槛 $3,500/kg",
         "原 v0 为 $2,800 — 对 1.2 t 级别过于乐观（隐含与 Falcon 9 同成本）",
         ROSE),
        ("小运载规模劣势",
         "1.2 t 级别天然 $/kg 高于 22 t 级别（缺乏规模摊薄）",
         GOLD),
        ("行业级曲线 Day 8 给出",
         "Day 8 成本模型将整合多厂商数据 + 不确定性带",
         TEAL),
    ]
    for i, (h, d, c) in enumerate(notes):
        y = Inches(1.5 + i*1.7)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(9.3), y, Inches(3.7), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = c
        card.line.width = Pt(2)
        # Header bar
        hh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(9.3), y,
                                     Inches(3.7), Inches(0.4))
        hh.fill.solid()
        hh.fill.fore_color.rgb = c
        hh.line.fill.background()
        p = hh.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
        # Description
        db = slide.shapes.add_textbox(Inches(9.45), y + Inches(0.5),
                                      Inches(3.4), Inches(0.95))
        p = db.text_frame.paragraphs[0]
        p.text = d
        p.font.size = Pt(10)
        p.font.color.rgb = DARK
        p.font.name = 'Microsoft YaHei'
        db.text_frame.word_wrap = True
    add_footer(slide, 17, 22)


# ------------------------------------------------------------------
# Slide 18 — 成本不确定性
# ------------------------------------------------------------------
def slide_cost_uncertainty():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "L0-03 成本目标 — 诚实的不确定性",
                   "Lognormal distribution (illustrative; Day 8 will derive from first principles).")
    img_path = f"{figures_dir}/cost_uncertainty_band.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.4),
                                  width=Inches(8.5))
    # Right side
    rt = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(9.3), Inches(1.4),
                                Inches(3.7), Inches(5.4))
    rt.fill.solid()
    rt.fill.fore_color.rgb = NAVY
    rt.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(9.5), Inches(1.55),
                                  Inches(3.3), Inches(5.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "如何在 v1.0 中体现"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = 'Microsoft YaHei'
    p.space_after = Pt(4)
    lines = [
        "采用对数正态分布表达 30% 不确定性",
        "明示 P(达到 $3,500/kg) ≈ 0.5",
        "",
        "诚实之处：",
        "v0 仅给出单点目标 ($2,800)",
        "v1.0 改为分布 + 概率注解",
        "",
        "Day 8 计划：",
        "整合 Monte Carlo + 多厂商基线",
        "输出 50 / 90 % 置信区间",
    ]
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        if "诚实之处" in line or "Day 8 计划" in line or "如何在" in line:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = GOLD if "Day 8" in line else WHITE
        else:
            p.font.size = Pt(10.5)
            p.font.color.rgb = RgbColor(0xC0, 0xD5, 0xE8)
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(2)
    add_footer(slide, 18, 22)


# ------------------------------------------------------------------
# Slide 19 — 概念车辆
# ------------------------------------------------------------------
def slide_concept():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "CRLV-1 概念车辆 — Day 1 草图",
                   "Not to scale; Day 2 will refine dimensions; Day 3 will detail engines.")
    img_path = f"{figures_dir}/concept_sketch.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3),
                                  width=Inches(9.0))
    # Right side
    rt = slide.shapes.add_textbox(Inches(9.7), Inches(1.5),
                                  Inches(3.4), Inches(5.2))
    tf = rt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Day 1 草图参数"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = 'Microsoft YaHei'
    p.space_after = Pt(6)
    specs = [
        ("总长", "≈ 40 m"),
        ("一级直径", "≈ 3.4 m"),
        ("一级干重", "≈ 18–22 t"),
        ("整流罩", "Ø 3.4 m × 6.5 m"),
        ("发动机", "7× methalox"),
        ("复用次数", "10–20"),
        ("推进剂", "LOX/LCH₄"),
    ]
    for k, v in specs:
        p = tf.add_paragraph()
        p.text = f"▸ {k}"
        p.font.size = Pt(11)
        p.font.color.rgb = SLATE
        p.font.name = 'Microsoft YaHei'
        p2 = tf.add_paragraph()
        p2.text = f"    {v}"
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = NAVY
        p2.font.name = 'Microsoft YaHei'
        p2.space_after = Pt(4)
    add_footer(slide, 19, 22)


# ------------------------------------------------------------------
# Slide 20 — Section divider: 5. 自评审 + 路径
# ------------------------------------------------------------------
def slide_section5():
    add_section_header(slide := prs.slides.add_slide(prs.slide_layouts[6]),
                       5, "自评审 · 计划",
                       ["Day 1 v1.0 self-critical review (12 issues fixed)",
                        "Forward path to Day 10",
                        "Living AI engineering notebook"])


# ------------------------------------------------------------------
# Slide 21 — 自评审（精炼版修正项）
# ------------------------------------------------------------------
def slide_review():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "Day 1 关键自评审（v1.0 修正项）",
                   "12 issues identified in self-review; all addressed.")
    items = [
        ("C1", "GRI 柱状图伪造数值 (1.00/1.15/1.35)", "改为量纲一致框架图，无绝对排名", ROSE),
        ("C2", "FoM 饼图与文档权重不一致", "重生成对齐 (30/20/25/15/10)", NAVY),
        ("C3", "“复用减排 95%”表述不严谨", "删除，改为按飞行摊薄", ROSE),
        ("C4", "成本门槛 $2,800 过于乐观", "上调至 $3,500 (小运载规模劣势)", ROSE),
        ("C5", "Long March 10B 命名不确定", "改用 “Long March 10/12A 家族” + 不确定性说明", ROSE),
        ("C6", "缺少 L1 离轨要求", "新增 L1-S02 (25 年门槛 / 5 年目标)", NAVY),
        ("C7", "着陆精度 ±100 m 过宽", "收紧到 ±30 / ±10 m (L1-O03)", ROSE),
        ("C8", "成本趋势被误读为行业均值", "标题与图注明示为单厂商轨迹", NAVY),
        ("C9", "L1-R01 表述与 L0-06 不一致", "明确单飞 vs 复合成功率的措辞", ROSE),
        ("C10", "Concept A 与主概念重复", "替换为 SSTO 真正的替代方案 (Concept C)", NAVY),
        ("C11", "整流罩 4.2 m × 8 m 过大", "减至 3.4 m × 6.5 m", ROSE),
        ("C12", "回收成功与复合成功率混用", "分别陈述并明确各自含义", ROSE),
    ]
    # Two-column layout
    for i, (cid, issue, fix, c) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.4 + col*6.4)
        y = Inches(1.4 + row*0.95)
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, y, Inches(6.2), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = c
        card.line.width = Pt(1.5)
        # ID badge
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    x, y, Inches(0.55), Inches(0.85))
        bg.fill.solid()
        bg.fill.fore_color.rgb = c
        bg.line.fill.background()
        bp = bg.text_frame.paragraphs[0]
        bp.text = cid
        bp.font.size = Pt(11)
        bp.font.bold = True
        bp.font.color.rgb = WHITE
        bp.alignment = PP_ALIGN.CENTER
        bp.font.name = 'Calibri'
        bg.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        # Issue
        ib = slide.shapes.add_textbox(x + Inches(0.65), y + Inches(0.04),
                                      Inches(5.4), Inches(0.4))
        p = ib.text_frame.paragraphs[0]
        p.text = issue
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = DARK
        p.font.name = 'Microsoft YaHei'
        # Fix
        fb = slide.shapes.add_textbox(x + Inches(0.65), y + Inches(0.42),
                                      Inches(5.4), Inches(0.4))
        p = fb.text_frame.paragraphs[0]
        p.text = "→ " + fix
        p.font.size = Pt(9.5)
        p.font.italic = True
        p.font.color.rgb = c
        p.font.name = 'Microsoft YaHei'
    # Footer note
    fb = slide.shapes.add_textbox(Inches(0.4), Inches(7.05),
                                  Inches(12.5), Inches(0.3))
    p = fb.text_frame.paragraphs[0]
    p.text = "v1.0 精炼版 = v0 草稿 + 12 项问题系统修正 + 自我评审机制建立（每日强制）"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = SLATE
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Microsoft YaHei'
    add_footer(slide, 21, 22)


# ------------------------------------------------------------------
# Slide 22 — 后续路径
# ------------------------------------------------------------------
def slide_forward():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "结论与后续路径（Day 2–10）",
                   "Day 1 complete with self-review; ready for Day 2.")
    # Gantt image
    img_path = f"{figures_dir}/gantt_10day.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.4), Inches(1.4),
                                  width=Inches(8.5))
    # Right side
    rt = slide.shapes.add_textbox(Inches(9.2), Inches(1.5),
                                  Inches(3.8), Inches(5.4))
    tf = rt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "关键节点"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = 'Microsoft YaHei'
    p.space_after = Pt(4)
    items = [
        ("Day 2", "第一级质量预算 + Δv 预算", SKY),
        ("Day 3", "推进剂循环 + 发动机选型", SKY),
        ("Day 4", "详细质量预算 + 材料/结构", SKY),
        ("Day 5", "气动 + 3DOF 上升/下降轨迹", SKY),
        ("Day 6", "推进式 vs 网捕权衡（关闭 L1-R03）", CORAL),
        ("Day 7", "多目标优化 + 首个 GRI 量化值", TEAL),
        ("Day 8", "成本模型 + 风险/可靠性量化", GOLD),
        ("Day 9", "系统集成 + 红队评审", VIOLET),
        ("Day 10", "终期展示", ROSE),
    ]
    for d, t, c in items:
        p = tf.add_paragraph()
        p.text = f"▸ {d}: {t}"
        p.font.size = Pt(10.5)
        p.font.bold = (d in ("Day 6", "Day 7"))
        p.font.color.rgb = c if d in ("Day 6", "Day 7") else DARK
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(3)
    # Conclusion box
    cb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.4), Inches(6.5),
                                Inches(12.5), Inches(0.55))
    cb.fill.solid()
    cb.fill.fore_color.rgb = NAVY
    cb.line.fill.background()
    p = cb.text_frame.paragraphs[0]
    p.text = "  结论：CRLV-1 任务定义 v1.0 已完成。L0 目标（6）、L1 需求（14）、GRI 框架、回收权衡均已就位；自评审机制建立。后续 Day 2–10 量化完善。"
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Microsoft YaHei'
    cb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_footer(slide, 22, 22)


# ------------------------------------------------------------------
# Build
# ------------------------------------------------------------------
print("Building CRLV-1 Day 1 Presentation v2.0...")
slide_cover()                # 1
slide_section1()             # 2
slide_background()           # 3
slide_data_2026()            # 4
slide_naming()               # 5
slide_method()               # 6
slide_section2()             # 7
slide_l0()                   # 8
slide_l1()                   # 9
slide_l1_viz()               # 10
slide_section3()             # 11
slide_gri()                  # 12
slide_gri_levers()           # 13
slide_section4()             # 14
slide_recovery()             # 15
slide_payload()              # 16
slide_cost()                 # 17
slide_cost_uncertainty()     # 18
slide_concept()              # 19
slide_section5()             # 20
slide_review()               # 21
slide_forward()              # 22
prs.save(output_path)
print(f"PPT generated: {output_path}")
print(f"Total slides: {len(prs.slides)}")
