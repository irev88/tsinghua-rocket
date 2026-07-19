#!/usr/bin/env python3
"""
Generate Comprehensive 10-Day PPT Presentation
Merges PDF program flow with GitHub repository's detailed plans and deliverables.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os
import json

# ============ CONFIGURATION ============
OUTPUT_PATH = "/home/user/tsinghua-rocket/docs/final/10_Day_Complete_Presentation.pptx"
FIGURES_DIR = "/home/user/tsinghua-rocket/figures"
DAY01_FIGURES = f"{FIGURES_DIR}/day01"

# Professional color scheme
PRIMARY = RgbColor(0x1A, 0x3C, 0x6E)      # Deep navy
ACCENT = RgbColor(0x00, 0x7A, 0xCC)       # Bright blue
SECONDARY = RgbColor(0x00, 0xB8, 0xA9)    # Teal
ACCENT2 = RgbColor(0xE6, 0x7E, 0x22)      # Orange
DARK = RgbColor(0x2D, 0x34, 0x36)
LIGHT_BG = RgbColor(0xF8, 0xF9, 0xFA)
WHITE = RgbColor(0xFF, 0xFF, 0xFF)
GRAY = RgbColor(0x95, 0xA5, 0xA6)
DARK_BLUE = RgbColor(0x0D, 0x1B, 0x2A)
GREEN = RgbColor(0x27, 0xAE, 0x60)
RED_ACCENT = RgbColor(0xC0, 0x39, 0x2B)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

# ============ SLIDE BUILDING BLOCKS ============
def add_bg(slide, color=PRIMARY):
    prs = slide.part.package.presentation_part.presentation
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_header_bar(slide, color=PRIMARY):
    prs = slide.part.package.presentation_part.presentation
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    return header

def add_title_text(slide, text, left=Inches(0.5), top=Inches(0.25), width=Inches(12.3), height=Inches(0.7), size=Pt(26), bold=True, color=WHITE, alignment=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return box

def add_body_text(slide, items, left=Inches(0.5), top=Inches(1.4), width=Inches(12.3), height=Inches(5.5), size=Pt(15), color=DARK, bullet=True, line_spacing=Pt(8)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ("• " if bullet else "") + item
        p.font.size = size
        p.font.color.rgb = color
        p.space_after = line_spacing
    return box

def add_two_column(slide, left_items, right_items, left_title="", right_title="", left_width=Inches(5.8), gap=Inches(0.5)):
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), left_width, Inches(5.3))
    tf = left_box.text_frame
    tf.word_wrap = True
    if left_title:
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.space_after = Pt(6)
    for i, item in enumerate(left_items):
        p = tf.add_paragraph() if left_title or i > 0 else tf.paragraphs[0]
        p.text = "• " + item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)

    right_x = Inches(0.5) + left_width + gap
    right_box = slide.shapes.add_textbox(right_x, Inches(1.5), Inches(12.3) - left_width - gap - Inches(0.5), Inches(5.3))
    tf = right_box.text_frame
    tf.word_wrap = True
    if right_title:
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.space_after = Pt(6)
    for i, item in enumerate(right_items):
        p = tf.add_paragraph() if right_title or i > 0 else tf.paragraphs[0]
        p.text = "• " + item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)

def add_figure(slide, image_path, left=Inches(1.0), top=Inches(1.3), width=Inches(11.3), caption=""):
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width=width)
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

def add_stat_cards(slide, stats, title):
    add_header_bar(slide)
    add_title_text(slide, title)
    colors = [ACCENT, SECONDARY, ACCENT2, GREEN, RED_ACCENT, PRIMARY]
    for i, (label, value) in enumerate(stats):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.1)
        y = Inches(1.5 + row * 2.5)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = colors[i % len(colors)]
        card.line.width = Pt(3)
        val_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(1.0))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = colors[i % len(colors)]
        p.alignment = PP_ALIGN.CENTER
        lbl_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.2), Inches(3.4), Inches(0.7))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

def add_table_slide(slide, title, headers, rows, col_widths=None):
    add_header_bar(slide)
    add_title_text(slide, title)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = DARK
                paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(0xF0, 0xF4, 0xF8)

def add_code_snippet(slide, title, code_lines, left=Inches(0.5), top=Inches(1.4), width=Inches(12.3), height=Inches(5.5)):
    add_header_bar(slide)
    add_title_text(slide, title)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.name = 'Consolas'
        p.font.color.rgb = DARK
        p.space_after = Pt(2)

# ============ SLIDE CONTENT DEFINITIONS ============
def build_slides(prs):
    # ============================================================
    # SLIDE 1: TITLE SLIDE
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, DARK_BLUE)
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    add_title_text(slide, "AI Co-Design of a Reusable Rocket", Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0), Pt(44), True, WHITE, PP_ALIGN.CENTER)
    add_title_text(slide, "10-Day Summer Program — Complete Technical Journey", Inches(1.5), Inches(3.1), Inches(10.3), Inches(0.8), Pt(24), False, RgbColor(0xCC, 0xE5, 0xFF), PP_ALIGN.CENTER)
    add_title_text(slide, "Tsinghua University Summer Program  ·  July 2026", Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.6), Pt(18), False, RgbColor(0xAA, 0xCC, 0xEE), PP_ALIGN.CENTER)
    add_title_text(slide, "Conceptual Vehicle: CRLV-1  |  Payload: 1.2–2.0 t to LEO/SSO  |  Reusable First Stage", Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.6), Pt(16), False, SECONDARY, PP_ALIGN.CENTER)
    add_title_text(slide, "AI-Assisted Multidisciplinary Engineering Workflow", Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.5), Pt(14), False, GRAY, PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 2: PROGRAM OVERVIEW (from PDF)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide)
    add_title_text(slide, "Program Overview — 10 Days, One Conceptual RLV")
    add_body_text(slide, [
        "Objective: Design a conceptual reusable launch vehicle (RLV) using AI-assisted engineering workflows",
        "Focus: Requirements decomposition, MDO, trade-off analysis, uncertainty quantification, engineering communication",
        "Philosophy: Scientific process > perfect rocket  |  Novelty > replication  |  Documentation > results",
        "Team: 3–5 students + LLM copilot + Prof. Xu (Chief Engineer) + Graduate mentors",
        "Deliverables: Mission report, mass budgets, propulsion trade, trajectory sims, recovery concept, optimization, cost/risk, technical review, final presentation",
        "AI Protocol: Every major decision logged with prompt, response, rationale, and impact (see ai_logs/)"
    ])

    # ============================================================
    # SLIDE 3: DAILY SCHEDULE (from PDF)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide)
    add_title_text(slide, "10-Day Schedule (from Program PDF)")
    headers = ["Day", "Date (2026)", "Theme", "Primary Deliverable", "Key Focus"]
    rows = [
        ["1", "Jul 11", "Mission Definition", "Mission Requirements", "Stakeholders, CONOPS, L0/L1 reqs, sustainability"],
        ["2", "Jul 12", "Rocket Fundamentals", "First-Order Sizing", "Rocket equation, mass fractions, strawman vehicle"],
        ["3", "Jul 13", "Propulsion System", "Engine Selection Trade", "Engine DB, methalox vs kerolox, Pugh matrix"],
        ["4", "Jul 14", "Mass Budget & Materials", "Vehicle Architecture", "Subsystem mass roll-up, CFRP/Al-Li/AM, TPS"],
        ["5", "Jul 15", "Aerodynamics & Trajectory", "Flight Profile Report", "3DOF ascent/descent, heating, landing dispersion"],
        ["6", "Jul 16", "Reusability Strategy", "Recovery Concept", "Propulsive vs net-capture, turnaround, GNC"],
        ["7", "Jul 17", "AI-Assisted Optimization", "Design Iteration", "MDO, Pareto front, LLM mutation, PINN surrogate"],
        ["8", "Jul 18", "Reliability & Economics", "Cost & Risk Analysis", "Cost model, Monte Carlo, business case, GRI"],
        ["9", "Jul 19", "System Integration", "Technical Review", "Mass closure, interface control, red-team critique"],
        ["10", "Jul 20", "Design Competition", "Final Presentation", "Story arc, demo, notebook export, repo package"],
    ]
    add_table_slide(slide, "", headers, rows, [Inches(0.6), Inches(1.0), Inches(2.2), Inches(3.5), Inches(5.0)])

    # ============================================================
    # SLIDE 4: GUIDING PRINCIPLES (from 10_DAY_BLUEPRINT)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide)
    add_title_text(slide, "Guiding Principles (from 10_DAY_BLUEPRINT.md)")
    left_items = [
        "Scientific Process First: Hypotheses → models → validation → iteration → documented rationale",
        "AI as Copilot, Not Oracle: Ideation, code, literature synthesis, trade-offs — always log prompts + decisions",
        "Novelty Preference: Bio-inspired recovery, PINN surrogates, sustainable propellants, explainable MOO, metamaterial TPS",
        "Academic Professionalism: Reproducible artifacts, UQ, balanced trade-offs, clear communication",
        "Scope Realism: Conceptual vehicle, first-order models + selective high-fidelity pockets, 1–2 novel contributions",
    ]
    right_items = [
        "Cross-Cutting Threads (every day):",
        "  • AI interaction logging (ai_logs/)",
        "  • Uncertainty & sensitivity analysis",
        "  • Sustainability / lifecycle angle",
        "  • Benchmarking vs real vehicles (Falcon 9, Zhuque-3, Neutron, LM-10B)",
        "  • Engineering notebook updates",
        "",
        "Tools Stack:",
        "  • Python (NumPy, SciPy, Matplotlib, Pandas, SymPy)",
        "  • RocketPy / OpenRocket / custom 3DOF",
        "  • PyTorch/JAX for surrogates, pymoo for MDO",
        "  • Plotly for interactive trades",
        "  • Markdown + LaTeX for reports",
    ]
    add_two_column(slide, left_items, right_items, "Core Principles", "Daily Practice & Tools")

    # ============================================================
    # SLIDE 5: DAY 1 — MISSION DEFINITION (Actual Deliverables)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, ACCENT)
    add_title_text(slide, "Day 1: Mission Definition — Completed Deliverables", color=WHITE)
    left_items = [
        "L0 Objectives (6): Payload 1.2–2.0 t, ≥10 reuse, <$2.8k/kg, CO₂e <15 t/t, 30-day responsiveness, 0.95 reliability",
        "L1 Requirements (15+): Performance, reusability, cost, environment, reliability, operations",
        "Novel: Green Reusability Index (GRI) as top-level FoM",
        "Novel: Explicit net-capture recovery evaluation (LM-10B 2026 demo)",
        "Stakeholder matrix, CONOPS, risk register, traceability skeleton",
        "Vehicle comparison table: 10 reference vehicles (2026 data)",
    ]
    right_items = [
        "Artifacts in docs/01_mission_requirements/:",
        "  • mission_requirements.md (v0.2, 92 lines)",
        "  • research_summary_2026.md (115 lines, 2025-26 data)",
        "  • vehicle_comparison_table.md (10 vehicles)",
        "  • stakeholder_analysis.md, initial_conops.md",
        "  • risk_register.md, requirements_traceability.md",
        "  • mission_concepts_alternatives.md",
        "",
        "AI Logs: prompts/day01_mission_definition.md, decisions/day01_initial.md",
        "Engineering notebook: Day 1 entry with hypotheses, decisions, open questions",
    ]
    add_two_column(slide, left_items, right_items, "Key Requirements & Innovations", "Deliverable Artifacts")

    # ============================================================
    # SLIDE 6: DAY 1 — KEY DATA & FIGURES
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, ACCENT)
    add_title_text(slide, "Day 1: Key Reference Data (2025–2026 Verified)", color=WHITE)
    stats = [
        ("Falcon 9 Reusable LEO", "22,800 kg"),
        ("Zhuque-3 Recovered LEO", "~18,300 kg"),
        ("LM-10B Net Capture LEO", "~16,000 kg"),
        ("Hyperbola-3 Recovered", "8,500 kg"),
        ("Pallas-1 Baseline", "8,000 kg"),
        ("RP-1/LOX CO₂e/t payload", "~19 t"),
        ("CRLV-1 Target Payload", "1,200–2,000 kg"),
        ("CRLV-1 Reuse Target", "10–20 flights"),
        ("Target Recurring Cost", "<$2,800/kg"),
    ]
    add_stat_cards(slide, stats, "2026 Benchmark Data & CRLV-1 Targets")

    # ============================================================
    # SLIDE 7: DAY 1 — FIGURES (Payload Comparison, Cost Trend, GRI)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, ACCENT)
    add_title_text(slide, "Day 1: Generated Figures", color=WHITE)
    # Three figure placeholders - using actual figures if they exist
    figs = [
        (f"{DAY01_FIGURES}/payload_comparison.png", "Reusable LEO Payload Comparison (2026)"),
        (f"{DAY01_FIGURES}/cost_trend.png", "Launch Cost Trend Driven by Reusability"),
        (f"{DAY01_FIGURES}/gri_comparison.png", "Green Reusability Index (GRI) Concept"),
    ]
    for i, (path, cap) in enumerate(figs):
        col = i % 3
        x = Inches(0.4 + col * 4.3)
        y = Inches(1.3)
        if os.path.exists(path):
            slide.shapes.add_picture(path, x, y, width=Inches(4.0))
        cap_box = slide.shapes.add_textbox(x, Inches(5.6), Inches(4.0), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = cap
        p.font.size = Pt(10)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

    # ============================================================
    # SLIDE 8: DAY 2 — ROCKET FUNDAMENTALS & SIZING (Plan + Code)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SECONDARY)
    add_title_text(slide, "Day 2: Rocket Fundamentals & First-Order Sizing", color=WHITE)
    left_items = [
        "Objectives: Rocket equation, staging laws, mass/dimension estimation, baseline strawman",
        "Novel: Symbolic regression for historical sizing laws, adaptive structural ratio by recovery mode",
        "Activities: Sizing script (code/sizing/), sensitivity on Isp & structural coeff, mass breakdown charts",
        "AI Co-pilot: Alternative architectures (SSTO reusable, TSTO, air-launch hybrid)",
        "Outputs: vehicle_sizing_report.md, baseline_mass_budget.csv, sizing_*.png",
        "Key parameters: Payload 1.2–2.0 t, Δv ~3.8 km/s (reusable), Isp 335 s (SL methalox)",
    ]
    right_items = [
        "Starter Code: code/sizing/first_order_sizer.py",
        "  • Tsiolkovsky rocket equation",
        "  • Dry mass estimation via structural coefficient",
        "  • CLI args: payload, Δv, Isp",
        "  • Placeholder for staging iteration",
        "",
        "Planned Sensitivity Studies:",
        "  • Structural coefficient: 0.06–0.12",
        "  • Isp: 320–350 s (methalox range)",
        "  • Recovery propellant margin: 5–20%",
        "  • Payload growth: 1.2→2.0 t impact",
        "",
        "Benchmarks: Falcon 9, Zhuque-3, Electron, Vega-C",
    ]
    add_two_column(slide, left_items, right_items, "Plan & Objectives", "Code & Analysis Scope")

    # Code snippet for Day 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SECONDARY)
    add_title_text(slide, "Day 2: first_order_sizer.py — Starter Code", color=WHITE)
    code_lines = [
        "import numpy as np",
        "G0 = 9.80665",
        "",
        "def rocket_equation_dv(m0, mf, isp):",
        "    return isp * G0 * np.log(m0 / mf)",
        "",
        "def estimate_dry_mass(payload, prop_mass, struct_coeff=0.08):",
        "    # struct_coeff = dry / (dry + prop)",
        "    return struct_coeff * (payload + prop_mass) / (1 - struct_coeff)",
        "",
        "def size_vehicle(payload_kg, dv_target, isp_sl, struct_coeff=0.08):",
        "    # Iterative sizing with staging",
        "    # TODO: Add upper stage, recovery margins",
        "    pass",
        "",
        "# CLI interface for quick trades",
        "if __name__ == '__main__':",
        "    parser.add_argument('--payload', default=1200)",
        "    parser.add_argument('--dv', default=3800)",
        "    parser.add_argument('--isp', default=335)",
    ]
    add_code_snippet(slide, "code/sizing/first_order_sizer.py", code_lines)

    # ============================================================
    # SLIDE 9: DAY 3 — PROPULSION SYSTEM
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, ACCENT2)
    add_title_text(slide, "Day 3: Propulsion System — Engine Selection Trade", color=WHITE)
    left_items = [
        "Objectives: Survey candidate engines, propulsion trade study, tank/feed sizing",
        "Novel: Green propellants (LOX/LCH4 vs additives), throttleable aerospike, dual-mode boost+recovery",
        "Physics-informed surrogate for chamber performance",
        "AI Co-pilot: Weighted Pugh matrix for 6+ candidate engines",
        "  Criteria: Isp, T/W, throttle range, reusability heritage, cost, manufacturability",
        "Outputs: engine_selection_trade.md, engines.csv, propulsion_trade.py",
    ]
    right_items = [
        "Candidate Engine Database (data/benchmarks/engines.csv):",
        "  • SpaceX Raptor 2/3 (methalox, 350s vac, deep throttle)",
        "  • Blue Origin BE-4 (methalox, 339s vac, 55% throttle)",
        "  • Landspace Tianque-12B (methalox, Zhuque-3, 9 engines)",
        "  • iSpace JD-1 (kerolox → methalox transition)",
        "  • Galactic Energy CQ-50 (kerolox, Pallas-1)",
        "  • CASC Longyun (methalox, LM-12A)",
        "  • Aerojet MRL10 (upper stage RL10-class)",
        "",
        "Trade Dimensions:",
        "  • Stage count (7–9 booster engines)",
        "  • Thrust-to-weight at liftoff (>1.3)",
        "  • Restart capability for recovery burns",
    ]
    add_two_column(slide, left_items, right_items, "Plan & Novel Angles", "Engine Candidates & Trade Space")

    # ============================================================
    # SLIDE 10: DAY 4 — MASS BUDGET & ADVANCED MATERIALS
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, GREEN)
    add_title_text(slide, "Day 4: Mass Budget & Advanced Materials — Vehicle Architecture", color=WHITE)
    left_items = [
        "Objectives: Detailed subsystem mass breakdown, material selection, structural sizing, tank design",
        "Novel Focus:",
        "  • 3D-printed additively manufactured tank domes + lattice structures",
        "  • Metamaterials / functionally graded materials for TPS",
        "  • Lifecycle mass: refurbishment mass penalty",
        "Activities: Mass budget with margins, material property DB, simple structural model, architecture diagram",
        "Outputs: vehicle_architecture.md, materials_db.json, structures/ subsystem notes",
    ]
    right_items = [
        "Subsystem Mass Breakdown Categories:",
        "  • Structures: Tanks, interstage, thrust structure, grid fins",
        "  • Propulsion: Engines, TVC, feed lines, pressurization",
        "  • Recovery: Legs (or net-capture interface), TPS, RCS",
        "  • Avionics: GNC, comm, power, harness",
        "  • Residuals & Margin: 10–15% system margin",
        "",
        "Material Trade Matrix:",
        "  • Al-Li 2195 (heritage, weldable, good cryo)",
        "  • CFRP (high specific strength, complex joints, cost)",
        "  • Ti alloys (high temp, recovery hardware)",
        "  • Stainless 304L (Starship approach, low cost, heavy)",
        "  • AM lattice cores (weight savings, inspection challenge)",
    ]
    add_two_column(slide, left_items, right_items, "Objectives & Novel Materials", "Mass Categories & Material Options")

    # ============================================================
    # SLIDE 11: DAY 5 — AERODYNAMICS & TRAJECTORY
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, PRIMARY)
    add_title_text(slide, "Day 5: Aerodynamics & Trajectory — Flight Profile", color=WHITE)
    left_items = [
        "Objectives: Conceptual aero config, 3DOF/6DOF sim, ascent + descent profiles, q_max, heating, landing footprint",
        "Novel Ideas:",
        "  • PINN (Physics-Informed NN) as fast trajectory surrogate",
        "  • Bio-inspired variable camber / morphing surfaces for recovery",
        "  • Monte Carlo dispersion on wind/guidance errors",
        "Activities: Implement/adapt trajectory code, generate altitude-velocity, heat rate, g-load plots",
        "AI Co-pilot: 3DOF simulator code + 3 non-standard recovery trajectories",
        "Outputs: flight_profile_report.md, trajectory_*.png, trajectory_sensitivity/",
    ]
    right_items = [
        "Key Trajectory Phases to Model:",
        "  1. Ascent: Pitch program, gravity turn, staging",
        "  2. Booster separation & boostback/entry burn",
        "  3. Reentry: Hypersonic → transonic, grid fin control",
        "  4. Landing burn: Terminal guidance, hover, touchdown",
        "",
        "Key Outputs for Trade Studies:",
        "  • Δv budget breakdown (ascent + recovery)",
        "  • Max q, max g, stagnation point heating",
        "  • Landing ellipse (3σ) vs propellant margin",
        "  • Grid fin sizing for control authority",
        "",
        "Tools: Custom 3DOF (Python), RocketPy for validation",
    ]
    add_two_column(slide, left_items, right_items, "Plan & Novel Approaches", "Trajectory Phases & Outputs")

    # ============================================================
    # SLIDE 12: DAY 6 — REUSABILITY STRATEGY (with brainstormed ideas)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, RED_ACCENT)
    add_title_text(slide, "Day 6: Reusability Strategy — Recovery Concept", color=WHITE)
    left_items = [
        "Objectives: Choose primary recovery method, size hardware, define ground ops & turnaround, assess reuse number",
        "Baseline: Propulsive landing + grid fins + legs (Falcon-like)",
        "Novel Alternative (from LM-10B 2026): Net capture on sea platform — no landing legs",
        "Brainstormed Novel Concepts (draft/brainstorm/novel_recovery_ideas.md):",
        "  1. Hybrid drone-assisted mid-air capture + retro",
        "  2. Bio-inspired variable geometry / morphing surfaces",
        "  3. Ocean platform with active wave compensation + magnetic capture",
        "  4. ML landing controller with propellant slosh damping",
        "  5. Circular recovery hardware (AM remanufacturing)",
        "Activities: Recovery mass model, landing accuracy vs propellant trade, turnaround timeline",
        "Outputs: recovery_concept.md, subsystems/recovery/, Day 6 trade study",
    ]
    right_items = [
        "Trade Study Dimensions:",
        "  • Mass penalty (legs ~1–2% dry mass vs net ship ops)",
        "  • Reliability (leg deployment vs net capture success)",
        "  • Turnaround time (pad vs ship return)",
        "  • Infrastructure cost (drone ship vs net vessel)",
        "  • Environmental impact (transport emissions)",
        "",
        "Key Metrics:",
        "  • Recovery success rate ≥90% after 5 flights",
        "  • Turnaround ≤30 days (goal 14)",
        "  • Refurb cost <10% of stage manufacturing",
        "  • Landing accuracy ≤100 m (3σ)",
        "",
        "AI Prompt: \"Compare 5 unconventional recovery architectures with rough mass penalty and TRL\"",
    ]
    add_two_column(slide, left_items, right_items, "Recovery Options & Brainstormed Concepts", "Trade Dimensions & Metrics")

    # ============================================================
    # SLIDE 13: DAY 7 — AI-ASSISTED OPTIMIZATION
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, RgbColor(0x8E, 0x44, 0xAD))  # Purple
    add_title_text(slide, "Day 7: AI-Assisted Optimization — Design Iteration", color=WHITE)
    left_items = [
        "Objectives: Multi-objective optimization (mass, cost, reliability, environment), major design iteration",
        "Novel Focus (from draft/brainstorm/optimization_novelty.md):",
        "  • LLM-guided evolutionary operators (mutation = architecture suggestion)",
        "  • PINN as fast surrogate inside optimizer",
        "  • Multi-fidelity Bayesian optimization",
        "  • Human-in-the-loop interactive Pareto steering via natural language",
        "AI Multi-Agent Session: Structures, Propulsion, Aero agents debate 5 radical variants",
        "Activities: Implement optimization/, run 2–3 cases, document Pareto front + selected point",
        "Outputs: optimization_report.md, updated baseline, strong notebook updates",
    ]
    right_items = [
        "Optimization Problem Formulation:",
        "  Objectives (weighted from Day 1 FoM):",
        "    1. Maximize reusable payload (30%)",
        "    2. Minimize recurring $/kg (25%)",
        "    3. Maximize reusability score (20%)",
        "    4. Maximize GRI (15%)",
        "    5. Minimize responsiveness days (10%)",
        "",
        "Design Variables:",
        "  • Booster diameter, length, engine count",
        "  • Structural coefficient, propellant margin",
        "  • Recovery architecture (legs vs net vs hybrid)",
        "  • Material choices per subsystem",
        "  • Upper stage config",
        "",
        "Code: code/optimization/multi_objective.py (pymoo)",
    ]
    add_two_column(slide, left_items, right_items, "Objectives & Novel AI Methods", "Optimization Formulation")

    # ============================================================
    # SLIDE 14: DAY 8 — RELIABILITY & ECONOMICS
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, RgbColor(0x2C, 0x3E, 0x50))
    add_title_text(slide, "Day 8: Reliability & Economics — Cost and Risk Analysis", color=WHITE)
    left_items = [
        "Objectives: First-order cost model, reliability block diagram, Monte Carlo cost/schedule, business case",
        "Novel Angles:",
        "  • Carbon pricing integrated into cost model",
        "  • \"Insurance\" model for reusability risk",
        "  • Sensitivity to flight rate & refurb time",
        "  • GRI as economic differentiator",
        "Activities: Cost model script, risk register + mitigations, sensitivity tornado plots",
        "Outputs: cost_risk_analysis.md, analysis/uq/, Monte Carlo results",
    ]
    right_items = [
        "Cost Model Components:",
        "  Development (non-recurring):",
        "    • Engine, structures, avionics, test, certification",
        "  Per-Flight (recurring):",
        "    • Propellant, refurbishment, recovery ops, range, insurance",
        "  Refurbishment: Labor, parts, inspection, recertification",
        "",
        "Reliability Approach:",
        "  • Fault tree / reliability block diagram",
        "  • Historical failure modes (engine, TPS, GNC, recovery)",
        "  • Bayesian updating with flight data",
        "",
        "Key Sensitivities:",
        "  • Flight rate (2→12/yr) impact on $/kg",
        "  • Refurb time (14→60 days) on fleet size",
        "  • Carbon price ($50–200/t CO₂e) on propellant choice",
    ]
    add_two_column(slide, left_items, right_items, "Objectives & Novel Economics", "Cost Model & Reliability")

    # ============================================================
    # SLIDE 15: DAY 9 — SYSTEM INTEGRATION
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, RgbColor(0x16, 0xA0, 0x85))
    add_title_text(slide, "Day 9: Final System Integration — Technical Review", color=WHITE)
    left_items = [
        "Objectives: Consolidate subsystems, integrated checks, identify open issues, prepare technical review",
        "Integration Checks:",
        "  • Mass closure (dry + prop + payload = liftoff)",
        "  • Performance closure (Δv budget with margins)",
        "  • Margin review (mass, performance, thermal, control)",
        "  • Interface compatibility (mech, electrical, fluid)",
        "AI Red-Team Review: \"List 10 strongest criticisms and potential fixes\"",
        "Activities: System-level summary, light ICD, red-team exercise, review package",
        "Outputs: technical_review_package/, updated engineering_notebook.md",
    ]
    right_items = [
        "Technical Review Package Contents:",
        "  1. Vehicle configuration drawing (block/exploded)",
        "  2. Mass properties report (CG, MoI, margins)",
        "  3. Performance summary (payload, orbit, Δv)",
        "  4. Key trade study decisions & rationale",
        "  5. Risk register with top 10 risks + mitigations",
        "  6. Cost estimate (dev + recurring + refurb)",
        "  7. GRI & sustainability assessment",
        "  8. Open issues & assumptions log",
        "",
        "Review Format:",
        "  • 20-min presentation + 15-min Q&A",
        "  • Peer-style critique (mentors + AI red team)",
        "  • Action items for Day 10 finalization",
    ]
    add_two_column(slide, left_items, right_items, "Integration Activities", "Review Package Contents")

    # ============================================================
    # SLIDE 16: DAY 10 — DESIGN COMPETITION & SHOWCASE
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, RgbColor(0xE7, 0x4C, 0x3C))
    add_title_text(slide, "Day 10: Design Competition & Showcase — Final Presentation", color=WHITE)
    left_items = [
        "Objectives: 10–15 min presentation, supporting visuals, interactive demo, final notebook, polished repo",
        "Presentation Story Arc:",
        "  1. Problem: Why this mission? (constellations, sustainability, responsiveness)",
        "  2. Process: Scientific method + AI copilot workflow",
        "  3. Concept: CRLV-1 configuration & key specs",
        "  4. Trade-offs: Key decisions (propellant, recovery, materials) with evidence",
        "  5. Novel Insight: GRI, net-capture evaluation, AI-driven iteration",
        "  6. Conclusions: Performance, cost, risk, path forward",
        "Activities: Rehearse narrative, generate one-pager + poster, record AI highlights",
        "Outputs: docs/final/, figures/final/, FINAL_PRESENTATION.md/.pptx, clean repo",
    ]
    right_items = [
        "Final Repository Package:",
        "  • All docs/ folders polished",
        "  • Code/ runnable with requirements.txt",
        "  • Data/ with READMEs for each dataset",
        "  • Figures/ with generation scripts",
        "  • ai_logs/ complete prompt/decision trail",
        "  • engineering_notebook.md (full journey)",
        "  • draft/ cleaned or archived",
        "",
        "Showcase Elements:",
        "  • Live trajectory plot (Plotly interactive)",
        "  • Pareto front explorer",
        "  • Mass budget waterfall chart",
        "  • GRI vs cost trade visualization",
        "  • AI decision log timeline",
        "",
        "Success Metrics (Academic Lens):",
        "  • Decision process transparency",
        "  • Evidence of iteration & learning",
        "  • At least one novel contribution",
        "  • Professional formatting & reproducibility",
    ]
    add_two_column(slide, left_items, right_items, "Presentation Narrative & Activities", "Final Package & Success Criteria")

    # ============================================================
    # SLIDE 17: CROSS-CUTTING THEMES — AI COPILOT WORKFLOW
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, RgbColor(0x34, 0x49, 0x5E))
    add_title_text(slide, "Cross-Cutting: AI Copilot Workflow (from draft/brainstorm/ai_collaboration_techniques.md)")
    left_items = [
        "Mandatory Protocol: Every major decision logged in ai_logs/ with prompt, model, response, decision, rationale",
        "Advanced Techniques Used:",
        "  1. Multi-Agent Role Play: Chief Engineer / Propulsion Lead / Structures Skeptic / Cost Accountant debate",
        "  2. Red Team / Blue Team: Attack current design with worst-case assumptions",
        "  3. Generative Design Mutation: LLM proposes new architecture vectors for evolutionary algorithm",
        "  4. Explainable Decision Extraction: AI writes technical review justifications",
        "  5. Prompt Chaining for UQ: List assumptions → assign uncertainties → qualitative Monte Carlo",
        "  6. Historical Failure Injection: \"List 5 historical RLV failures related to X. Redesign to avoid.\"",
    ]
    right_items = [
        "Logging Standard (ai_logs/):",
        "  • prompts/dayXX_topic.md — exact prompts used",
        "  • decisions/dayXX_decision_log.md — outcomes & rationale",
        "  • full_transcripts/ — optional raw responses",
        "",
        "Prompt Templates (draft/early_notes/ai_prompt_templates.md):",
        "  • General Engineering Copilot (role, state, task, requirements)",
        "  • Trade Study (weighted matrix, sensitivity, risks)",
        "  • Novelty Injection (unconventional approaches + feasibility)",
        "",
        "Stretch Goal: \"AI Engineering Copilot Best Practices\" appendix in final report",
    ]
    add_two_column(slide, left_items, right_items, "Advanced AI Collaboration Methods", "Logging Standards & Templates")

    # ============================================================
    # SLIDE 18: CROSS-CUTTING THEMES — SUSTAINABILITY & GRI
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, GREEN)
    add_title_text(slide, "Cross-Cutting: Sustainability & Green Reusability Index (GRI)", color=WHITE)
    left_items = [
        "From draft/brainstorm/sustainability_metrics.md:",
        "  GRI = Payload_kg / (CO₂e_kg_per_flight + Refurb_Energy_Penalty)",
        "",
        "Components Tracked Across All Days:",
        "  • Propellant production emissions (methane vs RP-1 vs H₂)",
        "  • Manufacturing emissions (CFRP energy intensity vs Al)",
        "  • Refurbishment energy + materials per flight",
        "  • Recovery transport emissions (ship vs truck vs drone)",
        "  • End-of-life disposal / recyclability",
        "",
        "Trades Enabled by GRI:",
        "  • LOX/LCH4 vs LOX/RP-1: methane wins on emissions but has boil-off",
        "  • Bio-derived RP or synthetic methane pathways",
        "  • Ground landing vs ocean (transport emissions)",
        "  • 3D-printed structures (lower buy-to-fly ratio)",
        "",
        "Optimization Target (Day 7): Maximize GRI alongside traditional FoMs",
    ]
    right_items = [
        "Data Sources to Investigate:",
        "  • LCA papers on Falcon 9, Ariane, Vega",
        "  • Propellant production LCA databases",
        "  • Carbon price sensitivity ($/t CO₂e scenarios)",
        "  • Upper atmosphere effects (H₂O, NOx, black carbon)",
        "",
        "Novel Contribution:",
        "  Instead of minimizing dry mass alone, optimize for",
        "  \"emissions per kg-to-orbit over 20 flights\"",
        "",
        "Early Hypothesis:",
        "  A slightly heavier vehicle using sustainable",
        "  propellants + shorter refurb time could beat",
        "  a lighter but dirtier design on GRI",
        "",
        "Alignment with 2025–2026 Trends:",
        "  • EU/China policy pushing LCA in design",
        "  • Methalox shift in all new reusable vehicles",
        "  • Reusability itself reduces manufacturing emissions/flight",
    ]
    add_two_column(slide, left_items, right_items, "GRI Formulation & Trade Space", "Data Sources & Novel Contribution")

    # ============================================================
    # SLIDE 19: REPOSITORY STRUCTURE & ARTIFACT MAP
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide)
    add_title_text(slide, "Repository Structure & Artifact Map (PROJECT_STRUCTURE.md)")
    # Use a tree-like text representation
    tree_text = [
        "tsinghua-rocket/",
        "├── 10_DAY_BLUEPRINT.md          # Day-by-day roadmap (living)",
        "├── summer_program.pdf           # Original program reference",
        "├── README.md                    # Project overview",
        "├── engineering_notebook.md      # Living design log (core artifact)",
        "├── PROJECT_STRUCTURE.md         # This structure doc",
        "├── docs/                        # Formal deliverables",
        "│   ├── 01_mission_requirements/ # Day 1 — 8 files completed",
        "│   ├── 02_rocket_fundamentals/  # Day 2",
        "│   ├── 03_propulsion/           # Day 3",
        "│   ├── 04_mass_budget_materials/# Day 4",
        "│   ├── 05_aerodynamics_trajectory/ # Day 5",
        "│   ├── 06_reusability_recovery/ # Day 6",
        "│   ├── 07_ai_optimization/      # Day 7",
        "│   ├── 08_reliability_economics/# Day 8",
        "│   ├── 09_system_integration/   # Day 9",
        "│   ├── 10_final_presentation/   # Day 10",
        "│   └── final/                   # Polished final package",
        "├── analysis/                    # Trade studies, sensitivity, UQ",
        "├── code/                        # Reusable Python scripts",
        "│   ├── sizing/                  # first_order_sizer.py",
        "│   ├── simulation/              # trajectory, aero",
        "│   ├── optimization/            # multi_objective.py (pymoo)",
        "│   └── utils/                   # helpers, constants",
        "├── simulations/                 # Simulation outputs & models",
        "├── data/                        # Inputs, outputs, benchmarks",
        "│   ├── inputs/mission_params.json",
        "│   ├── benchmarks/engines.csv",
        "│   └── outputs/*.csv",
        "├── figures/                     # All visuals (code + images)",
        "├── ai_logs/                     # AI copilot records (critical)",
        "│   ├── prompts/                 # Exact prompts",
        "│   └── decisions/               # Decision logs",
        "├── subsystems/                  # Deep dives per subsystem",
        "├── references/                  # Papers, datasheets, .bib",
        "└── draft/                       # Temporary (cleaned by Day 10)",
        "    ├── brainstorm/              # Novel ideas",
        "    ├── early_notes/             # Assumptions, strawman",
        "    └── work_plans/              # Daily micro-plans",
    ]
    add_body_text(slide, tree_text, size=Pt(10), bullet=False, line_spacing=Pt(2))

    # ============================================================
    # SLIDE 20: MISSION PARAMETERS (from data/inputs/mission_params.json)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, ACCENT)
    add_title_text(slide, "Current Baseline Parameters (data/inputs/mission_params.json)", color=WHITE)
    with open("/home/user/tsinghua-rocket/data/inputs/mission_params.json") as f:
        params = json.load(f)
    
    left_items = [
        f"Vehicle Name: {params['mission']['name']}",
        f"Payload Threshold: {params['mission']['payload_threshold_kg']} kg",
        f"Payload Goal: {params['mission']['payload_goal_kg']} kg",
        f"Target Orbit: {params['mission']['target_orbit']}",
        f"Reuse Target: {params['mission']['reusability_target_flights']} flights",
        f"Reuse Goal: {params['mission']['reusability_goal_flights']} flights",
        f"Turnaround Target: {params['mission']['turnaround_days_target']} days",
        f"Turnaround Goal: {params['mission']['turnaround_goal_days']} days",
    ]
    right_items = [
        f"Max G-load: {params['constraints']['max_g_load_ascent']}",
        f"Max q: {params['constraints']['max_dynamic_pressure_kPa']} kPa",
        f"Landing Accuracy: {params['constraints']['landing_accuracy_m']} m",
        f"Refurb Cost Limit: {params['constraints']['first_stage_reuse_cost_percent']}%",
        "",
        "FoM Weights:",
        f"  Payload: {params['figures_of_merit_weights']['payload_reusable_kg']*100:.0f}%",
        f"  Cost/kg: {params['figures_of_merit_weights']['recurring_cost_per_kg_usd']*100:.0f}%",
        f"  Reusability: {params['figures_of_merit_weights']['reusability_score']*100:.0f}%",
        f"  GRI: {params['figures_of_merit_weights']['sustainability_gri']*100:.0f}%",
        f"  Responsiveness: {params['figures_of_merit_weights']['responsiveness_days']*100:.0f}%",
    ]
    add_two_column(slide, left_items, right_items, "Mission & Constraints", "Figures of Merit Weights")

    # Key assumptions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, ACCENT)
    add_title_text(slide, "Key Assumptions (from mission_params.json)", color=WHITE)
    assumptions = params['key_assumptions']
    items = [
        f"Staging: {assumptions['staging']}",
        f"Primary Propellant: {assumptions['primary_propellant']}",
        f"Recovery Baseline: {assumptions['recovery_baseline']}",
        f"Alternative Recovery: {assumptions['alternative_recovery']}",
        f"Sustainability Metric: {assumptions['sustainability_metric']}",
        "",
        "Reference Vehicles for Benchmarking:",
    ] + [f"  • {v}" for v in params['reference_vehicles']]
    add_body_text(slide, items, size=Pt(14))

    # ============================================================
    # SLIDE 21: ENGINEERING NOTEBOOK — PROCESS EVIDENCE
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, RgbColor(0x7F, 0x8C, 0x8D))
    add_title_text(slide, "Engineering Notebook — Evidence of Scientific Process", color=WHITE)
    left_items = [
        "Living document: engineering_notebook.md",
        "Day 1 Entry Includes:",
        "  • Hypotheses (3) with research backing",
        "  • Key verified insights (Falcon 9, Chinese vehicles, sustainability)",
        "  • Requirements work summary with artifacts",
        "  • Figures of Merit with weights",
        "  • AI interaction log references",
        "  • Decisions & rationale (4 major decisions)",
        "  • Open questions & high-uncertainty items",
        "  • Next steps for Day 2 transition",
        "",
        "Template for Days 2–10:",
        "  Focus | Hypotheses Tested | AI Prompts Used",
        "  Major Outputs | Decisions & Rationale",
        "  Open Questions | Next Steps",
    ]
    right_items = [
        "Purpose: Primary evidence of",
        "  • Scientific process adherence",
        "  • AI collaboration transparency",
        "  • Iteration & learning documentation",
        "  • Decision traceability",
        "",
        "Cross-References:",
        "  • ai_logs/prompts/ — exact prompts",
        "  • ai_logs/decisions/ — outcomes",
        "  • docs/ — formal deliverables",
        "  • code/ — reproducible models",
        "  • figures/ — generated visuals",
        "",
        "Academic Value:",
        "  Notebook > Perfect rocket",
        "  Process documentation is the core grade",
        "  Demonstrates engineering judgment",
    ]
    add_two_column(slide, left_items, right_items, "Notebook Structure & Day 1 Content", "Purpose & Academic Value")

    # ============================================================
    # SLIDE 22: NOVEL CONTRIBUTIONS SUMMARY
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, RgbColor(0x8E, 0x44, 0xAD))
    add_title_text(slide, "Targeted Novel Contributions (1–2 Core + Supporting)", color=WHITE)
    novel_items = [
        ("1. Green Reusability Index (GRI)", "Day 1–10", "Sustainability as top-level FoM with lifecycle CO₂e + refurb penalty; drives propellant, recovery, ops trades", "Initiated"),
        ("2. Net-Capture Recovery Evaluation", "Day 6", "Grounded in LM-10B July 2026 success; formal trade vs propulsive legs; mass & ops implications", "Initiated"),
        ("3. AI Multi-Agent Optimization", "Day 7", "LLM-guided evolutionary operators + PINN surrogate + human-in-the-loop Pareto steering", "Planned"),
        ("4. Carbon-Priced Cost Model", "Day 8", "Integrated carbon cost ($50–200/t) into recurring economics; sensitivity to propellant choice", "Planned"),
        ("5. Bio-Inspired Morphing Recovery Surfaces", "Day 6", "Variable geometry grid fins/control surfaces for ascent/descent optimization", "Planned"),
        ("6. Circular Recovery Hardware Design", "Day 6/10", "AM remanufacturing of landing hardware; design for disassembly & recycle", "Planned"),
    ]
    # Create table
    table_shape = slide.shapes.add_table(len(novel_items)+1, 4, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.0))
    table = table_shape.table
    headers = ["Novel Contribution", "Primary Day", "Description & Impact", "Status"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RgbColor(0x8E, 0x44, 0xAD)
    for i, (name, day, desc, status) in enumerate(novel_items):
        for j, val in enumerate([name, day, desc, status]):
            cell = table.cell(i+1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.LEFT if j == 2 else PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(0xF5, 0xEE, 0xF8)

    # ============================================================
    # SLIDE 23: RISK MITIGATION & CONTINGENCY
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, RED_ACCENT)
    add_title_text(slide, "Risk Mitigation & Contingency Plans (from 10_DAY_BLUEPRINT.md)", color=WHITE)
    left_items = [
        "Identified Risks:",
        "  • Time slippage: Prioritize core deliverables; use minimum viable models",
        "  • Model fidelity gaps: Always state assumptions & limitations explicitly",
        "  • AI hallucination: Cross-check numbers with literature / simple physics",
        "  • Scope creep: Freeze requirements after Day 2 (with change log)",
        "  • Integration failures: Daily mass/performance closure checks from Day 4",
        "",
        "Contingency Actions:",
        "  • If behind: Drop detailed UQ on Day 8 → qualitative + bounds",
        "  • If ahead: Add stretch novel analysis (PINN trajectory surrogate)",
        "  • If integration fails: Simplify architecture (fewer engines, fixed recovery)",
        "  • If AI unreliable: Fall back to conventional methods + document",
    ]
    right_items = [
        "Milestones & Checkpoints:",
        "  • End Day 3: Baseline vehicle exists (rough numbers)",
        "  • End Day 6: Reusable concept closed",
        "  • End Day 7: ≥1 major AI-driven iteration complete",
        "  • End Day 8: Quantified cost/risk",
        "  • End Day 9: All subsystems integrated, review ready",
        "  • End Day 10: Showcase ready",
        "",
        "Success Measurement (Academic Lens):",
        "  1. Quality & transparency of decision process",
        "  2. Evidence of iteration & learning",
        "  3. Clarity of communication",
        "  4. ≥1 novel or non-obvious insight",
        "  5. Professional formatting & reproducibility",
        "",
        "Remember: Journey & documented reasoning > perfect rocket",
    ]
    add_two_column(slide, left_items, right_items, "Risks & Mitigations", "Milestones & Success Criteria")

    # ============================================================
    # SLIDE 24: CLOSING / THANK YOU
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BLUE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    add_title_text(slide, "CRLV-1: AI Co-Design of a Reusable Rocket", Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0), Pt(40), True, WHITE, PP_ALIGN.CENTER)
    add_title_text(slide, "10-Day Journey Complete  —  From Mission Definition to Technical Review", Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.8), Pt(22), False, RgbColor(0xCC, 0xE5, 0xFF), PP_ALIGN.CENTER)
    add_title_text(slide, "Key Deliverables: Mission Requirements · Vehicle Sizing · Propulsion Trade · Mass Budget", Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.5), Pt(16), False, SECONDARY, PP_ALIGN.CENTER)
    add_title_text(slide, "Trajectory Sims · Recovery Concept · AI Optimization · Cost/Risk · Integrated Review", Inches(1.5), Inches(4.9), Inches(10.3), Inches(0.5), Pt(16), False, SECONDARY, PP_ALIGN.CENTER)
    add_title_text(slide, "Novel Contributions: GRI · Net-Capture Evaluation · AI Multi-Agent MDO · Carbon-Priced Economics", Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5), Pt(16), False, GREEN, PP_ALIGN.CENTER)
    add_title_text(slide, "清华大学暑期项目 · 2026年7月  ·  AI辅助工程设计工作流", Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.5), Pt(14), False, GRAY, PP_ALIGN.CENTER)

    return prs

# ============ MAIN ============
if __name__ == "__main__":
    print("Creating presentation...")
    prs = create_presentation()
    prs = build_slides(prs)
    prs.save(OUTPUT_PATH)
    print(f"PPT saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")